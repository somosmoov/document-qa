import streamlit as st
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
import mammoth
import requests
import cohere

# Inicialize o cliente Cohere
cohere_client = cohere.Client(st.secrets["api_cohere"])  # Substitua pela sua chave da API Cohere

# Função para ler arquivos PDF
def read_pdf(file):
    try:
        document = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in document:
            text += page.get_text()
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PDF: {e}")
        return ""

# Função para ler arquivos DOCX
def read_docx(file):
    try:
        document = Document(file)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOCX: {e}")
        return ""

# Função para ler arquivos DOC usando mammoth
def read_doc(file):
    try:
        result = mammoth.extract_raw_text(file)
        return result.value
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOC: {e}")
        return ""

# Função para ler arquivos PPT e PPTX
def read_ppt_pptx(file):
    try:
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PPT/PPTX: {e}")
        return ""

# Função para ler arquivos TXT e MD
def read_txt_md(file):
    try:
        return file.read().decode()
    except Exception as e:
        st.error(f"Erro ao ler arquivo TXT/MD: {e}")
        return ""

# Função para enviar o texto para a API da Hugging Face
def query_huggingface_api(document_text, question):
    api_url = "https://api-inference.huggingface.co/models/YOUR_MODEL"  # Substitua pelo seu modelo na Hugging Face
    headers = {
        "Authorization": "Bearer YOUR_HUGGINGFACE_API_KEY"  # Substitua pela sua chave da API
    }
    payload = {
        "inputs": f"{document_text}\n\nQuestion: {question}"
    }
    try:
        response = requests.post(api_url, headers=headers, json=payload)
        response.raise_for_status()
        return response.json()[0].get("generated_text", "Nenhuma resposta encontrada.")
    except requests.exceptions.RequestException as e:
        st.error(f"Erro ao consultar a API da Hugging Face: {e}")
        return "Erro ao consultar a API da Hugging Face."

# Função para enviar o texto para a API do Cohere
def query_cohere_api(document_text, question, model):
    prompt = f"{document_text}\n\nQuestion: {question}\nAnswer:"
    try:
        response = cohere_client.generate(
            model=model,  # Substitua pelo modelo desejado
            prompt=prompt,
            max_tokens=100,
            temperature=0.7
        )
        if response.generations and response.generations[0].text:
            return response.generations[0].text.strip()
        else:
            st.error("Resposta vazia ou malformada recebida da API do Cohere.")
            return "Erro: Resposta vazia ou malformada recebida da API do Cohere."
    except cohere.CohereAPIError as e:
        st.error(f"Erro na API do Cohere: {e.message}")
        return "Erro na API do Cohere."
    except requests.exceptions.RequestException as e:
        st.error(f"Erro de rede ao consultar a API do Cohere: {e}")
        return "Erro de rede ao consultar a API do Cohere."
    except Exception as e:
        st.error(f"Ocorreu um erro inesperado: {e}")
        return "Ocorreu um erro inesperado ao consultar a API do Cohere."

# Streamlit UI
st.title("📝 Carregue o Edital")

uploaded_file = st.file_uploader("Carregue o arquivo com o edital", type=("pdf", "docx", "doc", "ppt", "pptx", "txt", "md"))

question = st.text_input(
    "Faça um questionamento",
    placeholder="Pode fornecer um sumário?",
    disabled=not uploaded_file,
)

# Dropdown para selecionar o modelo
model = st.selectbox(
    "Selecione o modelo Cohere",
    ["xlarge", "large", "medium", "small"]  # Ajuste conforme os modelos disponíveis na sua conta Cohere
)

if uploaded_file and question:
    # Process the uploaded file based on its type
    if uploaded_file.type == "application/pdf":
        document_text = read_pdf(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        document_text = read_docx(uploaded_file)
    elif uploaded_file.type == "application/msword":
        document_text = read_doc(uploaded_file)
    elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        document_text = read_ppt_pptx(uploaded_file)
    else:
        document_text = read_txt_md(uploaded_file)

    if document_text:
        answer = query_cohere_api(document_text, question, model)
        st.write(answer)  # Exibe a resposta da API do Cohere
''' 
    if document_text:
        answer = query_huggingface_api(document_text, question)
        st.write(answer)  # Exibe a resposta da API da Hugging Face
'''        


# Certifique-se de instalar as dependências necessárias
# pip install pymupdf python-docx python-pptx mammoth requests

