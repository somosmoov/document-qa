import streamlit as st
from openai import OpenAI
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
import mammoth
import requests
import openpyxl 
from azure.identity import DefaultAzureCredential
from azure.keyvault.secrets import SecretClient
from langchain_community.document_loaders import PDFMinerLoader
from langchain_community.document_loaders import PyMuPDFLoader
import pymupdf4llm
import chardet


KEY_VAULT_URL = "https://keyvaultdesen.vault.azure.net/"

# Função para ler arquivos PDF
def read_pdf(file):
    try:
        md_text = pymupdf4llm.to_markdown(file)
        return md_text 
        '''
        document = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in document:
            text += page.get_text()
        return text
        '''
        '''#loader = PDFMinerLoader(file)
        loader = PyMuPDFLoader(file)
        docs = loader.load()
        return docs '''
    except Exception as e:
        st.error(f"Erro ao ler arquivo PDF: {e}")
        return ""

# Função para ler arquivos DOCX
def read_docx(file):
    try:
        document = Document(file)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOCX: {e}")
        return ""

# Função para ler arquivos DOC usando mammoth
def read_doc(file):
    try:
        result = mammoth.extract_raw_text(file)
        return result.value
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOC: {e}")
        return ""

# Função para ler arquivos PPT e PPTX
def read_ppt_pptx(file):
    try:
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PPT/PPTX: {e}")
        return ""

# Função para ler arquivos TXT e MD com detecção automática de codificação
def read_txt_md(file, encoding="utf-8"):
    try:
        # Detectar a codificação do arquivo
        detected_encoding = chardet.detect(file.read())["encoding"]

        # Reabrir o arquivo com a codificação detectada
        file.seek(0)
        text = file.read().decode(encoding or detected_encoding)

        return text

    except Exception as e:
        st.error(f"Erro ao ler arquivo TXT/MD: {e}")
        return ""

def read_xls(file_path):
    try:
        # Abrir o arquivo XLS
        workbook = openpyxl.load_workbook(file_path)

        # Acessar a primeira planilha do arquivo
        worksheet = workbook.active

        # Criar um dicionário para armazenar os dados
        data = {}

        # Iterar sobre as linhas da planilha
        for row in worksheet.iter_rows():
            # Iterar sobre as células da linha atual
            for cell in row:
                # Verificar se a célula contém dados
                if cell.value:
                    # Adicionar os dados ao dicionário
                    data[cell.coordinate] = cell.value

        # Fechar o arquivo
        workbook.close()

        return data

    except FileNotFoundError:
        #print(f"Arquivo '{file_path}' não encontrado.")
        st.write(f"Arquivo '{file_path}' não encontrado.")
        return None

    except Exception as e:
        #print(f"Erro ao ler o arquivo XLS: {e}")
        sr.write(f"Erro ao ler o arquivo XLS: {e}")
        return None

# Função para converter o arquido carregado para texto
def trata_arquivo (uploaded_file):
    #st.write("Tipo arquivo carregado",uploaded_file.type)
    # Process the uploaded file based on its type
    if uploaded_file.type == "application/pdf":
        document_text = read_pdf(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        document_text = read_docx(uploaded_file)
    elif uploaded_file.type == "application/msword":
        document_text = read_doc(uploaded_file)
    elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        document_text = read_ppt_pptx(uploaded_file)
    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","application/xls","application/xlsx", "application/xlsm","application/xltx","application/xltm"]: 
        document_text = read_xls(uploaded_file)
    else:
        document_text = read_txt_md(uploaded_file)
    return document_text

# Função para limpar o conteúdo da variável "question"
def clear_question():
    try:
        del question
    except NameError:
        pass

# Criar uma caixa de entrada de texto
def get_question():
    clear_question()
    return st.text_input(
        "Agora é só perguntar que a IA responde!",
        placeholder="Por exemplo: Pode fornecer um sumário?",
        disabled=not uploaded_file,
    )

# Streamlit UI
#st.title("📝 Carregue o Documento")
st.markdown("## 📝 Carregue o Documento")

# Ask user for their OpenAI API key via `st.text_input`.
# Alternatively, you can store the API key in `./.streamlit/secrets.toml` and access it
# via `st.secrets`, see https://docs.streamlit.io/develop/concepts/connections/secrets-management
#openai_api_key = st.text_input("OpenAI API Key", type="password")
#openai_api_key = st.secrets["api_openai"]

# Defina a URL do seu Key Vault
key_vault_url = KEY_VAULT_URL

# Crie um cliente para acessar o Key Vault
credential = DefaultAzureCredential()
client = SecretClient(vault_url=key_vault_url, credential=credential)

# Acesse o segredo
secret_name = "OpenAI-API-Key"
retrieved_secret = client.get_secret(secret_name)
#st.write("A chave da API da OpenAI é:", retrieved_secret.value)
#st.write("VALOR recuperado:", retrieved_secret)

# Create an OpenAI client.
client = OpenAI(api_key=retrieved_secret.value)

# Let the user upload a file via `st.file_uploader`.
uploaded_file = st.file_uploader("Carregue o arquivo com o documento a ser analisado!", type=("pdf", "docx", "doc", "ppt", "pptx", "txt", "md","xls","xlsx","xlsm","xltx","xltm"))

# Chamar a função para obter a pergunta
# question = get_question()

# Ask the user for a question via `st.text_area`.
question = st.text_input(
    "Pergunte que a IA responde!",
    placeholder="Por exemplo: Pode fornecer um sumário?",
    disabled=not uploaded_file,
)

if uploaded_file and question:

    # Process the uploaded file and question.
    #document = uploaded_file.read().decode()
    document = trata_arquivo(uploaded_file)
    messages = [
        {
            "role": "user",
            "content": f"Here's a document: {document} \n\n---\n\n {question}",
        }
    ]
    #st.write(document)
    # Generate an answer using the OpenAI API.
    stream = client.chat.completions.create(
        #model="gpt-3.5-turbo",
        model="gpt-4o-mini",
        messages=messages,
        stream=True,
    )

    # Stream the response to the app using `st.write_stream`.
    st.write_stream(stream)

